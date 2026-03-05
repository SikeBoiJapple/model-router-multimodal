from __future__ import annotations

import base64
import io
import uuid
from dataclasses import dataclass
from typing import Iterable

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


class FilePipelineError(ValueError):
    pass


FILE_KIND_BY_EXTENSION: dict[str, str] = {
    ".jpg": "image",
    ".jpeg": "image",
    ".png": "image",
    ".webp": "image",
    ".gif": "image",
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
}

MIME_BY_EXTENSION: dict[str, str] = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".webp": "image/webp",
    ".gif": "image/gif",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

MAX_FILE_BYTES_BY_KIND: dict[str, int] = {
    "image": 10 * 1024 * 1024,
    "pdf": 20 * 1024 * 1024,
    "docx": 15 * 1024 * 1024,
    "xlsx": 15 * 1024 * 1024,
    "pptx": 20 * 1024 * 1024,
}


@dataclass(frozen=True)
class ProcessedFile:
    file_id: str
    filename: str
    extension: str
    kind: str
    mime_type: str
    size_bytes: int
    b64_data: str | None = None
    extracted_text: str | None = None


class FileStore:
    def __init__(self) -> None:
        self._items: dict[str, ProcessedFile] = {}

    def put(self, item: ProcessedFile) -> None:
        self._items[item.file_id] = item

    def get(self, file_id: str) -> ProcessedFile:
        item = self._items.get(file_id)
        if item is None:
            raise FilePipelineError(f"Unknown file_id '{file_id}'")
        return item

    def get_many(self, file_ids: Iterable[str]) -> list[ProcessedFile]:
        return [self.get(file_id) for file_id in file_ids]

    def delete(self, file_id: str) -> ProcessedFile:
        item = self._items.pop(file_id, None)
        if item is None:
            raise FilePipelineError(f"Unknown file_id '{file_id}'")
        return item


def detect_extension(filename: str) -> str:
    dot = filename.rfind(".")
    return filename[dot:].lower() if dot >= 0 else ""


def process_uploaded_file(filename: str, content_type: str | None, data: bytes) -> ProcessedFile:
    if not filename:
        raise FilePipelineError("File name is required")
    extension = detect_extension(filename)
    kind = FILE_KIND_BY_EXTENSION.get(extension)
    if kind is None:
        valid = ", ".join(sorted(FILE_KIND_BY_EXTENSION.keys()))
        raise FilePipelineError(f"Unsupported file type '{extension}'. Allowed: {valid}")

    size_bytes = len(data)
    max_size = MAX_FILE_BYTES_BY_KIND[kind]
    if size_bytes > max_size:
        raise FilePipelineError(
            f"File '{filename}' exceeds max size for {kind} ({max_size // (1024 * 1024)}MB)"
        )

    mime_type = content_type or MIME_BY_EXTENSION[extension]
    file_id = str(uuid.uuid4())
    b64_data: str | None = None
    extracted_text: str | None = None

    if kind in {"image", "pdf"}:
        b64_data = base64.b64encode(data).decode("ascii")

    if kind == "pdf":
        extracted_text = _extract_pdf_text(data)
    elif kind == "docx":
        extracted_text = _extract_docx_text(data)
    elif kind == "xlsx":
        extracted_text = _extract_xlsx_text(data)
    elif kind == "pptx":
        extracted_text = _extract_pptx_text(data)

    return ProcessedFile(
        file_id=file_id,
        filename=filename,
        extension=extension,
        kind=kind,
        mime_type=mime_type,
        size_bytes=size_bytes,
        b64_data=b64_data,
        extracted_text=extracted_text,
    )


def _extract_pdf_text(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    chunks: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text:
            chunks.append(text)
    return "\n".join(chunks).strip()


def _extract_docx_text(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    chunks = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(chunks).strip()


def _extract_xlsx_text(data: bytes) -> str:
    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"[Sheet] {sheet.title}")
        row_count = 0
        for row in sheet.iter_rows(values_only=True):
            cells = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if not cells:
                continue
            lines.append(" | ".join(cells))
            row_count += 1
            if row_count >= 200:
                lines.append("[Truncated after 200 non-empty rows]")
                break
    return "\n".join(lines).strip()


def _extract_pptx_text(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"[Slide {idx}]")
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                lines.append(text.strip())
    return "\n".join(lines).strip()
